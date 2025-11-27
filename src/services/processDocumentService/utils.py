import hashlib
import os
from uuid import uuid4
from openai import OpenAI
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pinecone import Pinecone, ServerlessSpec
from datetime import datetime
import mimetypes
import docx
from dotenv import load_dotenv
import fitz
import re
import mimetypes
from pptx import Presentation
from pdf2image import convert_from_path
from PIL import Image
import pytesseract
import unicodedata
from PIL import Image, ImageOps

load_dotenv()

# pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
pytesseract.pytesseract.tesseract_cmd = os.getenv("TESSERACT_CMD", pytesseract.pytesseract.tesseract_cmd)

# CONFIGURACIÓN INICIAL
EMBEDDING_MODEL = "text-embedding-3-small"
INDEX_NAME = "ciudadano-digital"
BATCH_SIZE = 100

openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
pinecone_client = Pinecone(api_key=os.getenv("PINECONE_API_KEY"))


# INICIALIZA PICONE Y VERIFICA EXISTENCIA DEL ÍNDICE
def init_pinecone():
    api_key = os.getenv("PINECONE_API_KEY")
    if not api_key:
        raise ValueError("No se encontró la variable de entorno PINECONE_API_KEY")

    pc = Pinecone(api_key=api_key)
    existing_indexes = [idx.name for idx in pc.list_indexes()]
    if INDEX_NAME not in existing_indexes:
        pc.create_index(
            name=INDEX_NAME,
            dimension=1536,
            metric="cosine",
            spec=ServerlessSpec(cloud="aws", region="us-east-1")
        )
    return pc.Index(INDEX_NAME)


index = init_pinecone()


# FUNCIONES AUXILIARES
def classify_category(fragment: str, categories: list) -> str:
    """Clasifica automáticamente el fragmento en una categoría."""
    prompt = f"""
    Clasifica el siguiente texto en una de las categorías:
    {categories} o bien, identifica una nueva categoría (dame el nombre de la categoría, no me digas "Nueva categoría") en caso consideres que no aplica dentro de ninguna de las opciones.
    Tampoco me des categorías compuestas, fuerza al resultado a ser una ÚNICA categoría.
    Responde SOLO con el nombre de la categoría. No hagas categorías compuestas, dame una única categoría.

    Texto: {fragment[:150]}
    """
    try:
        response = openai_client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}]
        )
        category = response.choices[0].message.content.strip()
        return category
    except Exception:
        return "General"


def already_indexed(sha1_hash: str) -> bool:
    """Verifica si un hash ya está en el índice (indexación incremental)."""
    query = index.query(
        vector=[0.0] * 1536,  # vector vacío ficticio
        top_k=1,
        filter={"sha1": {"$eq": sha1_hash}}
    )
    return len(query.get("matches", [])) > 0


def segment_text(text: str):
    """Divide el texto en fragmentos semánticamente coherentes de 20–150 palabras. Usa títulos si existen, si no hace split por palabras."""    
    # Dividir por títulos marcados
    sections = re.split(r'(?=<TITLE>)', text)
    chunks = []
    
    for section in sections:
        section = section.replace("<TITLE>", "").strip()
        if not section:
            continue
        
        words = section.split()
        if len(words) <= 150:
            chunks.append(section)
        else:
            start = 0
            while start < len(words):
                end = min(start + 150, len(words))
                fragment = " ".join(words[start:end])
                if len(fragment.split()) > 50:
                    chunks.append(fragment.strip())
                start = end
    
    # Si no se encontraron títulos y solo un fragmento gigante
    if not chunks and len(text.split()) > 150:
        words = text.split()
        start = 0
        while start < len(words):
            end = min(start + 150, len(words))
            chunks.append(" ".join(words[start:end]))
            start = end

    return [c for c in chunks if len(c.split()) >= 20]


def clean_text(text: str) -> str:
    """Limpieza profunda: texto plano, sin caracteres especiales ni espacios extra."""
    
    # NFKD separa tildes de letras
    text = unicodedata.normalize("NFKD", text)
    
    # Reemplazar saltos de línea y tabulaciones por un solo espacio
    text = re.sub(r'[\r\n\t]+', ' ', text)
    # Mantiene letras, números, espacios y puntuación básica"- 
    text = re.sub(r"[^a-zA-Z0-9\s.,;:¡!¿?()'\"-]", "", text)
    # Reducir múltiples espacios a uno solo
    text = re.sub(r'\s+', ' ', text)
    # Eliminar espacios iniciales/finales
    return text.strip()


def standardize_format(text: str) -> str:
    """Convierte títulos y numeración a formatos consistentes. Marca títulos con <TITLE> para segmentación."""
    # Títulos estilo Markdown o encabezados numéricos
    text = re.sub(r'(?m)^(#+\s*)(.+)$', lambda m: m.group(1) + "<TITLE> " + m.group(2).upper(), text)
    text = re.sub(r'(?m)^(\d+[\)\.-]\s*)(.+)$', lambda m: m.group(1) + "<TITLE> " + m.group(2).title(), text)
    # Títulos completamente en mayúsculas
    text = re.sub(r'(?m)^([A-Z][A-Z\s]{3,})$', lambda m: "<TITLE> " + m.group(1).title(), text)

    # Uniforma numeraciones: “1)”, “1.-”, “1.” → “1.”
    text = re.sub(r'(?m)^(\d+)[\)\.-]\s*', r'\1. ', text)

    return text.strip()


def validate_integrity(text: str) -> str:
    """Elimina duplicados, líneas corruptas y asegura coherencia del texto."""
    lines = text.splitlines()
    clean_lines = []

    for line in lines:
        if len(line.strip()) < 3:
            continue
        # Si más del 60% de los caracteres son válidos (alfanuméricos o espacio), se conserva
        valid_chars = sum(c.isalnum() or c.isspace() for c in line)
        if valid_chars / max(len(line), 1) > 0.6:
            clean_lines.append(line)

    # Elimina duplicados
    seen = set()
    deduped_lines = []
    for line in clean_lines:
        if line not in seen:
            deduped_lines.append(line)
            seen.add(line)
    clean_lines = deduped_lines

    return "\n".join(clean_lines).strip()


def extract_text_from_file(file_path: str) -> str:
    """Extrae texto de archivos PDF, Word, PowerPoint, texto plano o imágenes (usando OCR)."""
    mime_type, _ = mimetypes.guess_type(file_path)
    text = ""

    # --- PDF ---
    if mime_type == "application/pdf":
        with fitz.open(file_path) as pdf:
            for page_number, page in enumerate(pdf, start=1):
                page_text = page.get_text()
                # Si no hay texto, intentamos OCR en la imagen de la página
                if not page_text.strip():
                    pix = page.get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    # Preprocesamiento de imagen para OCR
                    img = ImageOps.grayscale(img)
                    img = ImageOps.invert(img)
                    page_text = pytesseract.image_to_string(img)
                    
                text += page_text + "\n"
        return text.strip()

    # --- Documentos Word ---
    elif mime_type in [
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ]:
        doc = docx.Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    # --- Presentaciones PowerPoint ---
    elif mime_type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
        return text.strip()

    elif mime_type == "application/vnd.ms-powerpoint":
        raise ValueError("Archivos .ppt antiguos no están soportados. Convierte a .pptx primero.")

    # --- Imágenes (JPG, PNG, etc.) ---
    elif mime_type and mime_type.startswith("image/"):
        img = Image.open(file_path)
        # Preprocesamiento para OCR
        img = ImageOps.grayscale(img)
        img = ImageOps.invert(img)
        return pytesseract.image_to_string(img)

    # --- Archivos de texto ---
    elif mime_type and mime_type.startswith("text/"):
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    else:
        raise ValueError(f"Tipo de archivo no soportado: {file_path}")


# INDEXACIÓN DEL DOCUMENTO
def process_and_index_document(file_path: str, source_title: str, author: str, year: int, identifier: str, categories: list, minAge: int, maxAge: int):
    """Procesa e indexa un documento completo en Pinecone."""
    category = "General"

    text = extract_text_from_file(file_path)
    clean = clean_text(text)
    standardized = standardize_format(clean)
    validated = validate_integrity(standardized)
    fragments = segment_text(validated)

    batch = []
    for _, frag in enumerate(fragments):
        sha1_hash = hashlib.sha1(frag.encode("utf-8")).hexdigest()

        # Evita reindexar fragmentos duplicados
        if already_indexed(sha1_hash):
            # Fragmento ya indexado, omitido.
            continue
        
        if not frag.strip():
            # Fragmento vacío, omitido.
            continue

        category = classify_category(frag, categories)

        # Crear embedding
        emb = openai_client.embeddings.create(
            model=EMBEDDING_MODEL,
            input=frag
        ).data[0].embedding

        metadata = {
            "document_id": identifier,
            "text": frag,
            "source": source_title,
            "author": author,
            "year": year,
            "category": category,
            "sha1": sha1_hash,
            "uploaded_at": datetime.now().isoformat(),
            "minAge":minAge,
            "maxAge": maxAge
        }

        batch.append({"id": str(uuid4()), "values": emb, "metadata": metadata})

        if len(batch) >= BATCH_SIZE:
            index.upsert(vectors=batch, namespace="ciudadania")
            batch = []

    if batch:
        index.upsert(vectors=batch, namespace="ciudadania")
        
    return {
        "success": True,
        "category":category
    }

    
# ELIMINAR DOCUMENTO
def delete_document(identifier: str):
    """Elimina todos los fragmentos asociados a un documento en Pinecone."""
    try:
        index.delete(
            filter={"document_id": {"$eq": identifier}},
            namespace="ciudadania"
        )
        return {"success": True, "deleted_document": identifier, "error":None}
    except Exception as e:
        return {"success": False, "error": str(e), "deleted_document": None}
